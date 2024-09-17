import sys
from pptx import Presentation


def odp_to_text(source, target):
    with open(target, "w") as stream:
        prs = Presentation(source)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    stream.write(shape.text + "\n")


if __name__ == "__main__":
    odp_to_text(sys.argv[1], sys.argv[2])
